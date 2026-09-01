#!/usr/bin/env python3
"""Premium PDF export helpers for FileMint.

Targets:
  * XLSX: extract native PDF tables where possible, then fall back to grouped
    text rows so the workbook is still editable.
  * PPTX: preserve exact page visuals as slide backgrounds and add an
    invisible editable/searchable text layer when text can be extracted.
  * HTML: create a self-contained visual preview with selectable text spans.

This is a local/offline converter. For scanned PDFs, OCR is attempted when
Tesseract is available; otherwise the report clearly marks visual fallback.
"""

from __future__ import annotations

import argparse
import base64
import html
import json
import os
import tempfile
from typing import Any

import fitz

from pdf_export_model import TextLine, TextWord, clean_text, make_report, points_to_emu
from pdf_export_spreadsheet import (
    export_xlsx,
    native_tables,
    row_segments,
    safe_sheet_name,
    style_sheet,
    words_to_grid,
)
from pdf_export_text import (
    join_positioned_words,
    maybe_resolve_ocr,
    native_words,
    ocr_words_for_page,
    page_text_lines,
    page_text_words,
    render_page_png,
    update_text_metrics,
    words_to_lines,
)

__all__ = [
    "TextLine",
    "TextWord",
    "clean_text",
    "export_html",
    "export_pptx",
    "export_xlsx",
    "join_positioned_words",
    "main",
    "maybe_resolve_ocr",
    "native_tables",
    "native_words",
    "ocr_words_for_page",
    "page_text_lines",
    "page_text_words",
    "render_page_png",
    "row_segments",
    "safe_sheet_name",
    "set_pptx_run_transparent",
    "style_sheet",
    "update_text_metrics",
    "words_to_grid",
    "words_to_lines",
    "write_report",
]


def export_html(
    src: str, dst: str, report: dict[str, Any], lang: str, text_layer: bool
) -> None:
    doc = fitz.open(src)
    ocr_lang = maybe_resolve_ocr(src, lang, text_layer, report)
    parts = [
        "<!doctype html>",
        '<html lang="und">',
        "<head>",
        '<meta charset="utf-8">',
        '<meta name="viewport" content="width=device-width, initial-scale=1">',
        "<title>FileMint PDF Export</title>",
        "<style>",
        "body{margin:0;background:#2a2f35;font-family:Arial,sans-serif;}",
        ".page{position:relative;margin:18px auto;background:#fff;box-shadow:0 4px 24px rgba(0,0,0,.28);overflow:hidden;}",
        ".page img{position:absolute;inset:0;width:100%;height:100%;}",
        ".text{position:absolute;white-space:pre;color:transparent;line-height:1;transform-origin:left top;}",
        ".text::selection{background:rgba(45,121,255,.35);color:transparent;}",
        "</style>",
        "</head>",
        "<body>",
    ]
    with tempfile.TemporaryDirectory() as tmpdir:
        try:
            for page in doc:
                image = os.path.join(tmpdir, f"page-{page.number + 1}.png")
                px_w, px_h = render_page_png(page, image, dpi=160)
                with open(image, "rb") as image_file:
                    data = base64.b64encode(image_file.read()).decode("ascii")
                width = round(page.rect.width * 160 / 72.0, 2)
                height = round(page.rect.height * 160 / 72.0, 2)
                parts.append(
                    f'<section class="page" style="width:{width}px;height:{height}px">'
                )
                parts.append(
                    f'<img alt="Page {page.number + 1}" src="data:image/png;base64,{data}">'
                )
                if text_layer:
                    lines = page_text_lines(
                        page, tmpdir, ocr_lang, report, allow_ocr=True
                    )
                    update_text_metrics(report, lines)
                    scale = 160 / 72.0
                    for line in lines:
                        left = line.x0 * scale
                        top = line.y0 * scale
                        fs = max(6.0, line.font_size * scale)
                        safe = html.escape(line.text)
                        parts.append(
                            f'<span class="text" style="left:{left:.2f}px;top:{top:.2f}px;font-size:{fs:.2f}px">{safe}</span>'
                        )
                parts.append("</section>")
            report["hiddenTextLayer"] = bool(
                text_layer and report["editableTextDetected"]
            )
            if report["hiddenTextLayer"]:
                report["notes"].append(
                    "HTML keeps exact page visuals and overlays selectable transparent text."
                )
            else:
                report["nonEditableVisualFallback"] = True
                report["warnings"].append(
                    "No text layer was available for this HTML export."
                )
        finally:
            doc.close()
    parts.append("</body></html>")
    with open(dst, "w", encoding="utf-8") as f:
        f.write("\n".join(parts))


def set_pptx_run_transparent(run: Any) -> None:
    from pptx.oxml.ns import qn
    from pptx.oxml.xmlchemy import OxmlElement

    r_pr = run._r.get_or_add_rPr()
    for child in list(r_pr):
        if child.tag == qn("a:solidFill"):
            r_pr.remove(child)
    solid = OxmlElement("a:solidFill")
    color = OxmlElement("a:srgbClr")
    color.set("val", "000000")
    alpha = OxmlElement("a:alpha")
    alpha.set("val", "0")
    color.append(alpha)
    solid.append(color)
    r_pr.append(solid)


def export_pptx(
    src: str, dst: str, report: dict[str, Any], lang: str, text_layer: bool
) -> None:
    from pptx import Presentation
    from pptx.util import Pt

    doc = fitz.open(src)
    if len(doc) == 0:
        raise RuntimeError("PDF contains no pages.")
    prs = Presentation()
    prs.slide_width = points_to_emu(doc[0].rect.width)
    prs.slide_height = points_to_emu(doc[0].rect.height)
    blank = prs.slide_layouts[6]
    ocr_lang = maybe_resolve_ocr(src, lang, text_layer, report)

    with tempfile.TemporaryDirectory() as tmpdir:
        try:
            for page in doc:
                slide = prs.slides.add_slide(blank)
                image = os.path.join(tmpdir, f"slide-{page.number + 1}.png")
                render_page_png(page, image, dpi=180)
                slide.shapes.add_picture(
                    image, 0, 0, width=prs.slide_width, height=prs.slide_height
                )
                if text_layer:
                    lines = page_text_lines(
                        page, tmpdir, ocr_lang, report, allow_ocr=True
                    )
                    update_text_metrics(report, lines)
                    sx = prs.slide_width / page.rect.width
                    sy = prs.slide_height / page.rect.height
                    for line in lines:
                        if not line.text.strip():
                            continue
                        left = int(line.x0 * sx)
                        top = int(line.y0 * sy)
                        width = max(
                            points_to_emu(8), int(max(8.0, line.x1 - line.x0) * sx)
                        )
                        height = max(
                            points_to_emu(7), int(max(7.0, line.y1 - line.y0) * sy)
                        )
                        box = slide.shapes.add_textbox(left, top, width, height)
                        box.text_frame.margin_left = 0
                        box.text_frame.margin_right = 0
                        box.text_frame.margin_top = 0
                        box.text_frame.margin_bottom = 0
                        p = box.text_frame.paragraphs[0]
                        run = p.add_run()
                        run.text = line.text
                        run.font.size = Pt(max(4.0, min(28.0, line.font_size)))
                        set_pptx_run_transparent(run)
            if text_layer and report["editableTextDetected"]:
                report["hiddenTextLayer"] = True
                report["notes"].append(
                    "PowerPoint slides preserve exact PDF visuals as backgrounds with editable transparent text boxes on top."
                )
            else:
                report["nonEditableVisualFallback"] = True
                report["warnings"].append(
                    "PowerPoint export used page images because no reliable text layer was available."
                )
        finally:
            doc.close()
    prs.save(dst)


def write_report(path: str | None, report: dict[str, Any]) -> None:
    if not path:
        return
    with open(path, "w", encoding="utf-8") as f:
        json.dump(report, f, ensure_ascii=False, indent=2)


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--input", required=True)
    parser.add_argument("--output", required=True)
    parser.add_argument("--target", required=True, choices=["xlsx", "pptx", "html"])
    parser.add_argument("--lang", default="auto")
    parser.add_argument("--table-detection", default="true")
    parser.add_argument("--text-layer", default="true")
    parser.add_argument("--report")
    args = parser.parse_args()

    target = args.target.lower()
    table_detection = str(args.table_detection).lower() not in {
        "0",
        "false",
        "no",
        "off",
    }
    text_layer = str(args.text_layer).lower() not in {"0", "false", "no", "off"}
    report = make_report(args.input, target)
    report["tableDetectionEnabled"] = table_detection
    report["layoutPreservationEnabled"] = True

    try:
        if target == "xlsx":
            export_xlsx(args.input, args.output, report, args.lang, table_detection)
        elif target == "pptx":
            export_pptx(args.input, args.output, report, args.lang, text_layer)
        elif target == "html":
            export_html(args.input, args.output, report, args.lang, text_layer)
        else:
            raise RuntimeError(f"Unsupported target: {target}")
    except Exception as exc:
        report["warnings"].append(str(exc))
        write_report(args.report, report)
        raise
    write_report(args.report, report)


if __name__ == "__main__":
    main()
