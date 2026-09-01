"""End-to-end tests for HTML/PPTX exports and the PDF export CLI."""

from __future__ import annotations

import json
import os
import sys
import tempfile
import unittest
from unittest.mock import patch

import fitz
from pptx import Presentation

SERVER_DIR = os.path.dirname(os.path.dirname(__file__))
if SERVER_DIR not in sys.path:
    sys.path.insert(0, SERVER_DIR)

import pdf_export
from pdf_export import export_html, export_pptx, write_report


def report() -> dict:
    return {
        "pdfType": "native",
        "tablesDetected": 0,
        "editableTextBoxes": 0,
        "editableCharacters": 0,
        "editableTextDetected": False,
        "ocrTextCandidates": 0,
        "lowConfidenceOcrAreas": 0,
        "ocrPasses": [],
        "textCoverageEstimate": 0,
        "nonEditableVisualFallback": False,
        "hiddenTextLayer": False,
        "warnings": [],
        "notes": [],
    }


def save_pdf(path: str, text: str | None) -> None:
    document = fitz.open()
    page = document.new_page(width=220, height=100)
    if text:
        page.insert_text((12, 35), text, fontsize=11)
    document.save(path)
    document.close()


class PdfExportOutputTests(unittest.TestCase):
    def test_export_html_embeds_page_visual_and_selectable_text(self) -> None:
        with tempfile.TemporaryDirectory() as tempdir:
            source = os.path.join(tempdir, "source.pdf")
            output = os.path.join(tempdir, "output.html")
            save_pdf(source, "FileMint & Co")
            metrics = report()

            export_html(source, output, metrics, "auto", text_layer=True)

            with open(output, encoding="utf-8") as file:
                contents = file.read()
            self.assertIn("data:image/png;base64,", contents)
            self.assertIn("FileMint &amp; Co", contents)
            self.assertIn('class="text"', contents)
            self.assertTrue(metrics["hiddenTextLayer"])
            self.assertTrue(metrics["editableTextDetected"])
            self.assertIn("selectable transparent text", metrics["notes"][0])

    def test_export_html_reports_missing_text_layer(self) -> None:
        with tempfile.TemporaryDirectory() as tempdir:
            source = os.path.join(tempdir, "blank.pdf")
            output = os.path.join(tempdir, "blank.html")
            save_pdf(source, None)
            metrics = report()

            export_html(source, output, metrics, "auto", text_layer=False)

            self.assertFalse(metrics["hiddenTextLayer"])
            self.assertTrue(metrics["nonEditableVisualFallback"])
            self.assertIn("No text layer", metrics["warnings"][0])

    def test_export_pptx_preserves_visual_and_editable_text_layer(self) -> None:
        with tempfile.TemporaryDirectory() as tempdir:
            source = os.path.join(tempdir, "source.pdf")
            output = os.path.join(tempdir, "output.pptx")
            save_pdf(source, "Editable FileMint")
            metrics = report()

            export_pptx(source, output, metrics, "auto", text_layer=True)

            presentation = Presentation(output)
            self.assertEqual(len(presentation.slides), 1)
            self.assertGreaterEqual(len(presentation.slides[0].shapes), 2)
            text = " ".join(
                shape.text
                for shape in presentation.slides[0].shapes
                if getattr(shape, "has_text_frame", False)
            )
            self.assertIn("Editable FileMint", text)
            self.assertTrue(metrics["hiddenTextLayer"])
            self.assertIn("transparent text boxes", metrics["notes"][0])

    def test_write_report_handles_optional_path_and_unicode(self) -> None:
        write_report(None, {"unused": True})
        with tempfile.TemporaryDirectory() as tempdir:
            output = os.path.join(tempdir, "report.json")
            write_report(output, {"message": "Файл"})
            with open(output, encoding="utf-8") as file:
                self.assertEqual(json.load(file), {"message": "Файл"})

    @patch("pdf_export.write_report")
    @patch("pdf_export.export_xlsx")
    @patch("pdf_export.make_report")
    def test_main_parses_boolean_options_and_dispatches_xlsx(
        self, make_report, export_xlsx, write_report_mock
    ) -> None:
        metrics = report()
        make_report.return_value = metrics
        argv = [
            "pdf_export.py",
            "--input",
            "in.pdf",
            "--output",
            "out.xlsx",
            "--target",
            "xlsx",
            "--table-detection",
            "off",
            "--text-layer",
            "no",
            "--report",
            "report.json",
        ]

        with patch.object(sys, "argv", argv):
            pdf_export.main()

        export_xlsx.assert_called_once_with(
            "in.pdf", "out.xlsx", metrics, "auto", False
        )
        write_report_mock.assert_called_once_with("report.json", metrics)
        self.assertFalse(metrics["tableDetectionEnabled"])
        self.assertTrue(metrics["layoutPreservationEnabled"])

    @patch("pdf_export.write_report")
    @patch("pdf_export.export_html", side_effect=RuntimeError("conversion failed"))
    @patch("pdf_export.make_report")
    def test_main_records_report_before_reraising_export_failure(
        self, make_report, _export_html, write_report_mock
    ) -> None:
        metrics = report()
        make_report.return_value = metrics
        argv = [
            "pdf_export.py",
            "--input",
            "in.pdf",
            "--output",
            "out.html",
            "--target",
            "html",
            "--report",
            "report.json",
        ]

        with patch.object(sys, "argv", argv):
            with self.assertRaisesRegex(RuntimeError, "conversion failed"):
                pdf_export.main()

        self.assertIn("conversion failed", metrics["warnings"])
        write_report_mock.assert_called_once_with("report.json", metrics)


if __name__ == "__main__":
    unittest.main()
